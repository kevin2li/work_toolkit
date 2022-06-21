import collections.abc  # cnanot delete
import os
import os.path as osp
import traceback
from pathlib import Path

import pptx
import typer
from docx import Document
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml.ns import qn
from docx.shared import Cm as doc_cm
from docx.shared import Inches
from docx.shared import Pt as doc_pt
from loguru import logger
from pptx import Presentation
from pptx.util import Cm, Pt

app = typer.Typer()

@app.command(name="extract", help="提取pptx备注")
def extract(
    input_path: str = typer.Argument(...),
    output_path: str = typer.Option(None, "--output-path", "-o", help="output path"),
):
    """提取pptx备注并保存为docx文件

    Args:
        input_path (str, optional): 输入pptx文件路径.
        output_path (str, optional): 输出docx保存路径.
    """
    try:
        prs = Presentation(input_path)
        slides = prs.slides
        document = Document()
        for i, slide in enumerate(slides, 1):
            if slide.has_notes_slide:
                text  = slide.notes_slide.notes_text_frame.text.strip()
                if not text:
                    text = "无"
            else:
                text = "无"
            p = document.add_paragraph('')
            p.paragraph_format.space_before = doc_pt(5)
            p.paragraph_format.space_after = doc_pt(0)
            t = p.add_run(f"第{i}页")
            t.bold = True
            t.font.name = '微软雅黑'
            text_list = text.split("\n")
            for item in text_list:
                if item.strip():
                    p = document.add_paragraph('')
                    t = p.add_run(item)
                    t.font.name = 'Times New Roman'
                    t.font.size = doc_pt(11)
                    t.element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
                    p.paragraph_format.space_after = doc_pt(0)
                    p.alignment = WD_PARAGRAPH_ALIGNMENT.JUSTIFY
                    p.paragraph_format.first_line_indent = t.font.size * 2
        if output_path is None:
            p = Path(input_path)
            output_path = osp.join(os.getcwd(), f"{p.stem}-备注.docx")
        document.save(output_path)
    except:
        logger.error("提取备注时发生错误！")
        traceback.print_exc()

if __name__ == "__main__":
    app()
